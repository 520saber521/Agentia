"""Document & PPT viewer — convert uploaded files to previewable HTML."""

from __future__ import annotations

import base64
import logging
import re
from pathlib import Path
from typing import Any

logger = logging.getLogger("agenthub.services.viewer")

EXT_MIME_MAP = {
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt": "application/vnd.ms-powerpoint",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".txt": "text/plain",
    ".csv": "text/csv",
    ".log": "text/plain",
}

PREVIEWABLE_DOC_EXTS = {".md", ".markdown", ".txt", ".csv", ".log"}


def _pptx_to_html(file_path: Path) -> str:
    """Convert a PPTX file to a self-contained HTML slides viewer."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation(str(file_path))

    slides_html: list[str] = []
    for i, slide in enumerate(prs.slides):
        elements: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # Escape HTML entities
                    escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    # Preserve line breaks
                    escaped = escaped.replace("\n", "<br>")
                    left = shape.left / 914400 if shape.left else 0  # EMU to inches
                    top = shape.top / 914400 if shape.top else 0
                    width = shape.width / 914400 if shape.width else 8
                    height = shape.height / 914400 if shape.height else 1
                    font_size = 16
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size:
                                    font_size = run.font.size / 12700  # EMU to pt
                                    break
                    elements.append(
                        f'<div class="slide-element" style="'
                        f'left:{left * 100 / 10:.1f}%;top:{top * 100 / 7.5:.1f}%;'
                        f'width:{width * 100 / 10:.1f}%;'
                        f'font-size:{font_size}px;'
                        f'">'
                        f'<p>{escaped}</p></div>'
                    )
            if shape.has_table:
                table = shape.table
                rows_html = []
                for row in table.rows:
                    cells = "".join(
                        f"<td>{cell.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;') or '&nbsp;'}</td>"
                        for cell in row.cells
                    )
                    rows_html.append(f"<tr>{cells}</tr>")
                elements.append(
                    f'<div class="slide-table">'
                    f'<table>{"".join(rows_html)}</table></div>'
                )

        slides_html.append(
            f'<div class="slide" id="slide-{i}" style="display:{ "block" if i == 0 else "none"}">'
            f'<div class="slide-num">{i + 1}/{len(prs.slides)}</div>'
            f'{"".join(elements)}'
            f'</div>'
        )

    return _ppt_viewer_template(len(prs.slides), "".join(slides_html))


def _ppt_viewer_template(total: int, slides: str) -> str:
    return f'''<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>PPT 预览 · {total} 页</title>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:-apple-system,BlinkMacSystemFont,"PingFang SC","Microsoft YaHei",sans-serif;
  background:#0f1117;color:#e8eaed;display:flex;flex-direction:column;height:100vh}}
.toolbar{{display:flex;align-items:center;justify-content:center;gap:12px;padding:14px 20px;
  background:#1a1c24;border-bottom:1px solid #2d3039;user-select:none}}
.toolbar button{{border:0;border-radius:8px;padding:8px 16px;cursor:pointer;
  background:#2d3039;color:#c5c9d2;font-size:14px;transition:background .15s}}
.toolbar button:hover{{background:#3d4049}}
.toolbar button:disabled{{opacity:.35;cursor:default}}
.toolbar .page-info{{font-size:14px;color:#8b909c;min-width:80px;text-align:center}}
.slides{{flex:1;overflow:hidden;display:grid;place-items:center;padding:24px}}
.slide{{position:relative;width:min(960px,90vw);aspect-ratio:4/3;background:white;
  border-radius:10px;box-shadow:0 8px 40px rgba(0,0,0,.5);overflow:hidden;color:#1f2329}}
.slide-num{{position:absolute;bottom:10px;right:14px;font-size:11px;color:#8b909c;z-index:2}}
.slide-element{{position:absolute;padding:4px 8px;line-height:1.6;overflow:hidden}}
.slide-element p{{margin:0}}
.slide-table{{position:absolute;left:5%;top:10%;width:90%;overflow:auto}}
.slide-table table{{width:100%;border-collapse:collapse;font-size:14px}}
.slide-table td{{border:1px solid #dadce0;padding:8px 12px}}
@media(max-width:640px){{.toolbar{{padding:10px 12px;gap:8px}}.toolbar button{{padding:6px 12px;font-size:12px}}}}
</style>
</head>
<body>
<div class="toolbar">
<button onclick="prev()" id="btn-prev">◀ 上一页</button>
<span class="page-info"><span id="cur">1</span> / {total}</span>
<button onclick="next()" id="btn-next">下一页 ▶</button>
<button onclick="toggleFullscreen()" title="全屏">⛶</button>
</div>
<div class="slides">{slides}</div>
<script>
let cur=0,total={total};
function show(n){{document.getElementById('slide-'+cur).style.display='none';
cur=(n+total)%total;document.getElementById('slide-'+cur).style.display='block';
document.getElementById('cur').textContent=cur+1;
document.getElementById('btn-prev').disabled=cur===0;
document.getElementById('btn-next').disabled=cur===total-1;}}
function prev(){{if(cur>0)show(cur-1);}}
function next(){{if(cur<total-1)show(cur+1);}}
function toggleFullscreen(){{if(!document.fullscreenElement){{document.documentElement.requestFullscreen();}}
else{{document.exitFullscreen();}}}}
document.addEventListener('keydown',e=>{{if(e.key==='ArrowLeft')prev();if(e.key==='ArrowRight')next();}});
document.getElementById('btn-prev').disabled=true;
if(total<=1){{document.getElementById('btn-next').disabled=true;}}
</script>
</body>
</html>'''


def _markdown_to_html(text: str) -> str:
    """Convert Markdown text to styled HTML."""
    escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    # Restore HTML entities after escaping
    lines = escaped.split("\n")
    html_lines: list[str] = []
    in_code_block = False
    code_lines: list[str] = []
    code_lang = ""

    i = 0
    while i < len(lines):
        line = lines[i]

        if line.startswith("```"):
            if in_code_block:
                code = "\n".join(code_lines)
                html_lines.append(f'<pre><code class="language-{code_lang}">{code}</code></pre>')
                code_lines = []
                in_code_block = False
            else:
                in_code_block = True
                code_lang = line[3:].strip() or "plaintext"
                code_lines = []
            i += 1
            continue

        if in_code_block:
            code_lines.append(line)
            i += 1
            continue

        # Headers
        if line.startswith("###### "):
            html_lines.append(f"<h6>{line[7:]}</h6>")
        elif line.startswith("##### "):
            html_lines.append(f"<h5>{line[6:]}</h5>")
        elif line.startswith("#### "):
            html_lines.append(f"<h4>{line[5:]}</h4>")
        elif line.startswith("### "):
            html_lines.append(f"<h3>{line[4:]}</h3>")
        elif line.startswith("## "):
            html_lines.append(f"<h2>{line[3:]}</h2>")
        elif line.startswith("# "):
            html_lines.append(f"<h1>{line[2:]}</h1>")
        elif line.startswith("- ") or line.startswith("* "):
            html_lines.append(f"<li>{line[2:]}</li>")
        elif re.match(r"^\d+\. ", line):
            html_lines.append(f"<li>{re.sub(r'^\d+\. ', '', line)}</li>")
        elif line.startswith("> "):
            html_lines.append(f'<blockquote>{line[2:]}</blockquote>')
        elif line.startswith("---"):
            html_lines.append("<hr>")
        elif line.strip() == "":
            html_lines.append("<br>")
        else:
            # Inline: bold, italic, code, links
            processed = line
            processed = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", processed)
            processed = re.sub(r"\*(.+?)\*", r"<em>\1</em>", processed)
            processed = re.sub(r"`(.+?)`", r"<code>\1</code>", processed)
            processed = re.sub(r"\[(.+?)\]\((.+?)\)", r'<a href="\2" target="_blank">\1</a>', processed)
            html_lines.append(f"<p>{processed}</p>")

        i += 1

    if in_code_block:
        html_lines.append(f'<pre><code>{"\n".join(code_lines)}</code></pre>')

    return _doc_viewer_template("Markdown Preview", "".join(html_lines))


def _text_to_html(text: str, title: str = "Document Preview") -> str:
    """Convert plain text to styled HTML."""
    escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    body = "\n".join(f"<p>{line or '<br>'}</p>" for line in escaped.split("\n"))
    return _doc_viewer_template(title, body)


def _doc_viewer_template(title: str, body: str) -> str:
    return f'''<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{title}</title>
<style>
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:-apple-system,BlinkMacSystemFont,"PingFang SC","Microsoft YaHei",sans-serif;
  background:#1e1e2e;color:#cdd6f4;max-width:860px;margin:0 auto;padding:32px 24px;line-height:1.8}}
h1,h2,h3,h4,h5,h6{{color:#cba6f7;margin:20px 0 10px}}
h1{{font-size:28px;border-bottom:1px solid #45475a;padding-bottom:8px}}
h2{{font-size:22px}}h3{{font-size:18px}}
p{{margin:8px 0}}
pre{{background:#11111b;border-radius:10px;padding:16px;overflow:auto;margin:12px 0}}
code{{font-family:"Cascadia Mono",Consolas,monospace;font-size:13px;background:#313244;padding:2px 6px;border-radius:4px}}
pre code{{background:transparent;padding:0}}
blockquote{{border-left:3px solid #cba6f7;padding:8px 16px;margin:12px 0;color:#a6adc8;background:#31324433;border-radius:0 8px 8px 0}}
li{{margin-left:24px;list-style:disc}}
a{{color:#89b4fa}}hr{{border:0;border-top:1px solid #45475a;margin:16px 0}}
</style>
</head>
<body>{body}</body>
</html>'''


def render_preview_html(file_path: Path, mime_type: str) -> str | None:
    """Convert a file to previewable HTML. Returns None if not supported."""
    ext = file_path.suffix.lower()

    if ext == ".pptx":
        try:
            return _pptx_to_html(file_path)
        except Exception as exc:
            logger.warning("PPTX conversion failed: %s", exc)
            return _text_to_html(f"PPTX conversion error: {exc}", "PPT Preview Error")

    if ext in PREVIEWABLE_DOC_EXTS:
        try:
            text = file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            text = file_path.read_text(encoding="latin-1", errors="replace")
        if ext in (".md", ".markdown"):
            return _markdown_to_html(text)
        return _text_to_html(text, file_path.name)

    return None


def can_preview(ext: str | None, mime_type: str | None = None) -> bool:
    """Check if a file type can be previewed in the viewer."""
    ext = (ext or "").lower()
    if ext in (".pptx", ".ppt"):
        return True
    if ext in PREVIEWABLE_DOC_EXTS:
        return True
    if mime_type == "text/html":
        return True
    return False
